import os
import time
import requests
import tempfile
import subprocess
from enum import Enum
from io import BytesIO
from typing import ClassVar
from collections.abc import Callable

import pptx
import docx
import docx2txt
import pdfplumber
from bs4 import BeautifulSoup

from streamlit.runtime.uploaded_file_manager import UploadedFile

from streamlitextras.logger import log
from ai.connectors.TokenHelper import TokenHelper

from gptroles.appsettings import settings
from utils import methodclass


class DocxExtractionMethods(str, Enum):
    text: str = "text"
    clauses: str = "clauses"
    joined_paragraphs: str = "joined_paragraphs"
    labelled_paragraphs: str = "labelled_paragraphs"


# Previously was "joined_paragraphs"
default_docx_extraction_method = DocxExtractionMethods.text


@methodclass
class FileInputParser:
    supported_formats: ClassVar = [
        "doc",
        "docx",
        "pdf",
        "pptx",
        "txt",
        "log",
        "csv",
        "wav",
        "mp3",
        # "mp4",
    ]

    def extract_speech_text(cls, file: UploadedFile) -> str:
        # Upload the file
        upload_response = requests.post(
            "https://api.assemblyai.com/v2/upload", files={"file": file}
        )
        upload_url = upload_response.json()["upload_url"]

        # Request transcription
        headers = {"authorization": "c2db62b3ff83418497a0f8f86f34fc09"}
        json = {"audio_url": upload_url}
        transcript_request_response = requests.post(
            "https://api.assemblyai.com/v2/transcript", json=json, headers=headers
        )

        # Poll for the transcription to be completed
        transcript_id = transcript_request_response.json()["id"]
        while True:
            transcript_response = requests.get(
                f"https://api.assemblyai.com/v2/transcript/{transcript_id}",
                headers=headers,
            )
            status = transcript_response.json()["status"]
            if status == "completed":
                break
            if status == "failed":
                raise Exception("Transcription failed")
            time.sleep(5)

        # Return the transcript
        return transcript_response.json()["text"]

    def extract_docx_text(cls, file: UploadedFile) -> str:
        text = docx2txt.process(file)
        return text

    def extract_docx_joined_paragraphs(
        cls, file: UploadedFile, prefix: str = "\n"
    ) -> str:
        paragraphs = cls.extract_docx_paragraphs_texts(file)
        return prefix.join(paragraphs)

    def extract_docx_paragraphs_texts(
        cls, file: UploadedFile, exclude_empty: bool = True
    ) -> list[str]:
        paragraphs = cls.extract_docx_paragraphs(file)
        return [
            p.text.strip()
            for p in paragraphs
            if not exclude_empty or bool(p.text.strip())
        ]

    def extract_docx_paragraphs(cls, file: UploadedFile) -> list:
        try:
            doc = docx.Document(file)
        except Exception:
            raise ValueError("The uploaded file is not a valid .docx file.")

        return list(doc.paragraphs)

    def extract_docx_clauses(cls, file: UploadedFile) -> list:
        paragraphs = cls.extract_docx_paragraphs(file)
        others = []
        clauses = []
        out_text = ""

        for para in paragraphs:
            # print("X", para.style, para.text[:30])
            para_text = para.text.strip()
            if not para_text:
                continue

            if para_text.endswith(".") and len(para_text.split()) > 6:
                clauses.append(para_text)
                out_text += para_text + "\nC: "
            else:
                others.append(para_text)
                out_text += para_text + "\n"
        # Remove trailing delimiter
        out_text = out_text.strip("C: \n")
        return out_text

    def extract_csv_text(cls, uploaded_file: UploadedFile) -> str:
        import pandas as pd

        try:
            # read the file with pandas
            dataframe = pd.read_csv(uploaded_file)

            # convert the dataframe to string
            csv_text = dataframe.to_string()

            return csv_text
        except Exception as e:
            log.error(f"An error occurred while reading the CSV file: {e}")
            return None

    def extract_pdf_text(cls, file: UploadedFile) -> str:
        try:
            with BytesIO(file.getvalue()) as data:
                with pdfplumber.open(data) as pdf:
                    full_text = []
                    for page in pdf.pages:
                        full_text.append(page.extract_text())
        except Exception as e:
            raise ValueError("The uploaded file is not a valid .pdf file.") from e

        return " ".join(full_text)

    def extract_pptx_text(cls, pptx_file: UploadedFile) -> str:
        prs = pptx.Presentation(pptx_file)
        text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    text.append(paragraph.text)

        return "\n".join(text)

    def read_file_contents(cls, uploaded_file: UploadedFile) -> str | None:
        content = uploaded_file.read().decode("utf-8")
        return content


@methodclass
class FileInputConvertor:
    conversion_directory = os.path.join(settings.tmp_folder, "document_conversions/")

    def convert_doc_to_docx(cls, doc_file: UploadedFile) -> UploadedFile:
        # log.debug(f"Conversion directory: {cls.conversion_directory}")
        os.makedirs(cls.conversion_directory, exist_ok=True)
        with tempfile.NamedTemporaryFile(
            suffix=".doc",
            dir=cls.conversion_directory,
        ) as tmp:
            # log.debug(f"temp file: {tmp.name}")
            tmp.write(doc_file.getvalue())

            result = subprocess.run(
                [
                    f"libreoffice --headless --convert-to docx {tmp.name} --outdir {cls.conversion_directory}"
                ],
                cwd=cls.conversion_directory,
                shell=True,
                stdout=subprocess.PIPE,
            )
            converted_file_name = tmp.name.replace(".doc", ".docx")
            converted_file_path = os.path.join(
                cls.conversion_directory, converted_file_name
            )
            if result.returncode == 0 and os.path.exists(converted_file_path):
                doc_file.seek(0)
                doc_file.truncate(0)
                with open(converted_file_path, "rb") as f:
                    doc_file.write(f.read())

                doc_file.name = (
                    doc_file.name.replace(".docx", "").replace(".doc", "") + ".docx"
                )
                doc_file.type = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                os.remove(converted_file_path)
                # log.success(f"Converted {doc_file.name}")
            else:
                log.error("Problem converting doc file to docx format:")
                log.warning(result.stdout.decode("utf-8"))

        return doc_file


@methodclass
class InputParser(FileInputParser, FileInputConvertor):
    def scrape_url(cls, url: str) -> str:
        if not url.startswith("http://") or url.startswith("https://"):
            url = "https://" + url

        try:
            response = requests.get(url)
        except requests.exceptions.ConnectionError:
            raise ValueError(
                "Unable to download from URL, please check it is valid."
            ) from None

        if response.status_code != 200:
            raise ValueError(
                f"Error {response.status_code}: Unable to fetch content from the provided URL."
            )
        soup = BeautifulSoup(response.content, "html.parser")
        text = soup.get_text(separator=" ")
        return text

    def join_text_from_files(
        cls,
        files: list[UploadedFile],
        docx_extraction_method: DocxExtractionMethods = default_docx_extraction_method,
    ) -> tuple[str, list[UploadedFile]]:
        new_files = []
        files_texts = []
        for file in files:
            file_contents_text, new_file = cls.extract_text_from_file(
                file, docx_extraction_method=docx_extraction_method
            )
            file = new_file
            new_files.append(new_file)
            files_texts.append((file.name, file_contents_text))

        concat_file_text = ""
        for file_name, text in files_texts:
            concat_file_text += f"File `{file_name}` contents:\n{text}\n"

        return concat_file_text, new_files

    def extract_text_from_file(
        cls,
        uploaded_file: UploadedFile,
        docx_extraction_method: DocxExtractionMethods = default_docx_extraction_method,
    ) -> str:
        extraction_method_functions = {
            DocxExtractionMethods.labelled_paragraphs: lambda x: cls.extract_docx_joined_paragraphs(
                x, "\n¶: "
            ),
            DocxExtractionMethods.joined_paragraphs: cls.extract_docx_joined_paragraphs,
            DocxExtractionMethods.clauses: cls.extract_docx_clauses,
            DocxExtractionMethods.text: cls.extract_docx_text,
        }

        if uploaded_file.type == "application/pdf":
            input_text = cls.extract_pdf_text(uploaded_file)
        elif (
            uploaded_file.type
            == "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        ):
            input_text = cls.extract_pptx_text(uploaded_file)
        elif (
            uploaded_file.type
            == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        ):
            extraction_function = extraction_method_functions.get(
                docx_extraction_method,
                extraction_method_functions[default_docx_extraction_method],
            )
            input_text = extraction_function(uploaded_file)
        elif uploaded_file.type == "application/msword":
            uploaded_file = cls.convert_doc_to_docx(uploaded_file)
            return cls.extract_text_from_file(uploaded_file, docx_extraction_method)
        elif uploaded_file.type == "text/csv":
            input_text = cls.extract_csv_text(uploaded_file)
        elif uploaded_file.type in ["audio/wav", "audio/mp3", "audio/mp4"]:
            input_text = cls.extract_speech_text(uploaded_file)
        else:
            input_text = uploaded_file.getvalue().decode("UTF-8")
        return input_text, uploaded_file

    def chunk_text(cls, text: str, chunk_length: int | None = None) -> list[str]:
        # TODO Implement boundary/paragraph checking to avoid splitting inside lexical chunks
        if not chunk_length:
            chunk_length = 6000

        words = text.split()
        chunks = []
        current_chunk = []
        current_token_count = 0

        for word in words:
            word_token_count = TokenHelper.calculate_token_length(word)
            if current_token_count + word_token_count <= chunk_length:
                current_chunk.append(word)
                current_token_count += word_token_count
            else:
                chunks.append(" ".join(current_chunk))
                current_chunk = [word]
                current_token_count = word_token_count

        if current_chunk:
            chunks.append(" ".join(current_chunk))

        return chunks
